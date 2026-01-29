#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enhanced PowerPoint extraction script - Version 3.6 (patched)
✅ Adds real PPT table extraction (shape.has_table) so "Composition" tables are captured
✅ Adds parsing for your typical AVO layout: [3 column headers] + [% row] like "31 | 67,5 | 1,5"
✅ Keeps previous behavior as fallback (chemical symbols parsing in plain text)

PATCH (Important):
- Fix comments propagation & entity_full_text aggregation across slides by using a stable group key.
- For "Matière première", the stable key is based on the reference (ref xxxx xxx) when available.
"""

import os
import re
import json
import argparse
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any

from pptx import Presentation
from PIL import Image
import io


def clean_text(text: str) -> str:
    """Clean and normalize text content."""
    if not text:
        return ""
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _extract_nuance_from_text(text: str) -> Optional[str]:
    """Extract nuance like '477 00' or '52307'."""
    if not text:
        return None

    t = " ".join(str(text).split())

    # Prefer explicit "Nuance"
    m = re.search(r"\bNuance\s*[:\-]?\s*(\d{3})\s*(\d{2})\b", t, flags=re.IGNORECASE)
    if m:
        return f"{m.group(1)} {m.group(2)}"

    # Generic "523 07"
    m = re.search(r"\b(\d{3})\s+(\d{2})\b", t)
    if m:
        return f"{m.group(1)} {m.group(2)}"

    # Compact "52307"
    m = re.search(r"\b(\d{5})\b", t)
    if m:
        s = m.group(1)
        return f"{s[:3]} {s[3:]}"
    return None


def extract_detailed_comments(text: str) -> Optional[str]:
    """Extract detailed comments/description from a slide."""
    if not text:
        return None

    # Method 1: Extract everything after "Commentaires:"
    comments_match = re.search(
        r"Commentaires\s*:(.*?)(?:\n\s*(?:Composition|TABLE|$))",
        text,
        re.IGNORECASE | re.DOTALL,
    )
    if comments_match:
        comments = comments_match.group(1).strip()
        comments = clean_text(comments)
        if len(comments) > 50:
            return comments

    # Method 2: Look for long text blocks
    paragraphs = re.split(r"\n\s*\n", text)

    for para in paragraphs:
        para_clean = clean_text(para)

        if len(para_clean) < 50:
            continue

        if re.match(r"^(Graphite|Coke|Nuance|Grossissement|Échelle)", para_clean, re.IGNORECASE):
            if len(para_clean) < 150:
                continue

        return para_clean

    return None


# ------------------------------
# NEW: Real PPT table extraction
# ------------------------------
def extract_tables_from_slide(slide) -> List[List[List[str]]]:
    """
    Return all tables as: [table][row][cell_text]
    Each table is a list of rows, each row is a list of cell strings.
    """
    tables: List[List[List[str]]] = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            try:
                tbl = shape.table
                table_rows: List[List[str]] = []
                for r in tbl.rows:
                    row_cells = []
                    for c in r.cells:
                        row_cells.append((c.text or "").strip())
                    table_rows.append(row_cells)
                tables.append(table_rows)
            except Exception:
                # Some shapes can be weird; fail gracefully
                continue
    return tables


def linearize_tables(tables: List[List[List[str]]]) -> str:
    """Convert extracted tables to a stable text form appended into raw_text."""
    if not tables:
        return ""
    blocks = []
    for t_idx, t in enumerate(tables, start=1):
        lines = []
        for row in t:
            # keep empty cells but trim right empties
            row2 = list(row)
            while row2 and row2[-1] == "":
                row2.pop()
            if not row2:
                continue
            lines.append(" | ".join(row2))
        if lines:
            blocks.append(f"TABLE_PPT_{t_idx}:\n" + "\n".join(lines))
    return "\n\n".join(blocks).strip()


# -----------------------------------------
# NEW: Parsing for "header + percent" tables
# -----------------------------------------
_NUM_TOKEN = re.compile(r"^[<>]?\s*\d+(?:[.,]\d+)?\s*$")


def _is_number_like(s: str) -> bool:
    return bool(_NUM_TOKEN.match((s or "").strip()))


def parse_avo_composition_from_tables(tables: List[List[List[str]]]) -> Optional[Dict[str, Any]]:
    """
    Looks for a table that matches the typical format:
      Row 1: component names (>=2 columns)
      Row N: numeric percentages aligned with headers
    Example:
      Cuivre Electrolytique Léger | Black mix 033 90 | Mélange abrasif 4042
      31 | 67,5 | 1,5

    Returns:
      {
        "elements": [...headers...],
        "values": [...values...],
        "rows": [headers, values]
      }
    """
    for t in tables:
        # Normalize table rows: strip cells and drop fully empty rows
        rows = []
        for r in t:
            r2 = [(c or "").strip() for c in r]
            if any(c for c in r2):
                rows.append(r2)

        if len(rows) < 2:
            continue

        # Find a candidate "values row" (mostly numeric)
        for idx in range(1, len(rows)):
            value_row = rows[idx]
            numeric_cells = [c for c in value_row if c and _is_number_like(c)]
            if len(numeric_cells) < 2:
                continue

            header_row = rows[0]

            pairs: List[Tuple[str, str]] = []
            max_cols = min(len(header_row), len(value_row))
            for j in range(max_cols):
                h = (header_row[j] or "").strip()
                v = (value_row[j] or "").strip()
                if not h:
                    continue
                if not v or not _is_number_like(v):
                    continue
                pairs.append((h, v))

            if len(pairs) >= 2:
                elements = [h for h, _ in pairs]
                values = [v for _, v in pairs]
                return {"elements": elements, "values": values, "rows": [elements, values]}

    return None


# ------------------------------
# Existing fallback parser (kept)
# ------------------------------
def extract_composition_table(text: str) -> Optional[Dict]:
    """Fallback: Extract composition table from plain text when it contains chemical symbols."""
    element_pattern = r"\b([A-Z][a-z]?)\b"
    lines = text.split("\n")

    for i, line in enumerate(lines):
        elements = re.findall(element_pattern, line)
        if len(elements) >= 3:
            if i + 1 < len(lines):
                values_line = lines[i + 1]
                values = re.findall(r"[<>]?\s*\d+(?:[.,]\d+)?", values_line)

                if len(values) >= len(elements):
                    return {
                        "elements": elements,
                        "values": [v.strip() for v in values[: len(elements)]],
                        "rows": [elements, [v.strip() for v in values[: len(elements)]]],
                    }
    return None


def composition_table_to_text(comp_table: Optional[Dict]) -> str:
    """Convert composition_table to readable text."""
    if not comp_table:
        return ""
    rows = comp_table.get("rows") or []
    if not rows:
        return ""
    lines = []
    for r in rows:
        lines.append(" | ".join([str(c).strip() for c in r]))
    return "\n".join(lines).strip()


# ✅ PATCH: Stable group key
def group_key(meta: Dict) -> Optional[str]:
    """
    A stable key to group slides that belong to the same entity.
    Problem: entity_id can vary if product_name is truncated on image slides.
    Solution: for "Matière première", use reference when available.
    """
    if not meta:
        return None
    if meta.get("entity_type") == "Matière première":
        ref = (meta.get("reference") or "").strip()
        if ref:
            return f"MP|{ref}"
    return meta.get("entity_id")


def extract_metadata_from_slide(slide, slide_number: int) -> Dict:
    """Extract metadata from a PowerPoint slide."""
    metadata = {
        "slide_number": slide_number,
        "nuance": None,
        "grade": None,
        "product_name": None,
        "reference": None,
        "magnification": None,
        "scale": None,
        "comments": None,
        "description": None,
        "composition": {},
        "composition_table": None,
        "full_text": None,
        "has_images": False,
        "entity_type": None,
        "entity_id": None,
    }

    # Extract all text + detect images
    full_text_blocks: List[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            full_text_blocks.append(shape.text)
        if hasattr(shape, "image"):
            metadata["has_images"] = True

    # NEW: Extract real PPT tables and append them to raw_text
    tables = extract_tables_from_slide(slide)
    tables_text = linearize_tables(tables)
    if tables_text:
        full_text_blocks.append(tables_text)

    raw_text = "\n".join(full_text_blocks)
    metadata["full_text"] = clean_text(raw_text)

    # Product name patterns (matières premières)
    product_patterns = [
        r"Graphite\s+(Timrex|SFG|KS|BNL|Naturel|Artificiel)\s*[A-Z0-9]*",
        r"Graphite\s+(?:artificiel|naturel)\s+[A-Za-z0-9\s]+",
        r"Coke\s+[A-Z]{2,}",
        r"graphite\s+(?:artificiel|naturel)[\s\n]+[A-Za-z0-9\s]+",
    ]
    for pattern in product_patterns:
        match = re.search(pattern, raw_text, re.IGNORECASE)
        if match:
            product = re.sub(r"\s+", " ", match.group(0).strip())
            metadata["product_name"] = product
            break

    # Reference
    ref_patterns = [
        r"(?:ref\s*[:\-]?\s*)?(\d{4,}\s+\d{2,})",
        r"-\s*(\d{4,}\s+\d{2,})",
        r"–\s*(\d{4,}\s+\d{2,})",
    ]
    for pattern in ref_patterns:
        match = re.search(pattern, raw_text, re.IGNORECASE)
        if match:
            ref_numbers = match.group(1).strip()
            metadata["reference"] = f"ref {ref_numbers}"
            break

    # Magnification
    mag_patterns = [
        r"Grossissement\s*[:\-]?\s*[xX×]?\s*(\d+)",
        r"[xX×]\s*(\d+)",
        r"(\d+)\s*[xX×]",
    ]
    for pattern in mag_patterns:
        match = re.search(pattern, raw_text)
        if match:
            mag_value = match.group(1)
            metadata["magnification"] = f"x{mag_value}"
            break

    # Scale
    scale_patterns = [
        r"[ÉéEe]chelle\s*[:\-]?\s*(\d+\s*[µμ]?m)",
        r"(\d+\s*[µμ]m)",
    ]
    for pattern in scale_patterns:
        match = re.search(pattern, raw_text, re.IGNORECASE)
        if match:
            metadata["scale"] = match.group(1).strip()
            break

    # Nuance detection
    nu = _extract_nuance_from_text(raw_text)
    if nu:
        metadata["nuance"] = nu

    # Detailed comments
    detailed_comments = extract_detailed_comments(raw_text)
    if detailed_comments:
        metadata["comments"] = detailed_comments
        metadata["description"] = detailed_comments

    # ---------------------------------------------------------
    # Composition extraction (NEW order):
    # 1) Try parsing real PPT tables in AVO format
    # 2) Fallback to old chemical-symbol text parsing
    # ---------------------------------------------------------
    comp_table = parse_avo_composition_from_tables(tables)
    if not comp_table:
        comp_table = extract_composition_table(raw_text)

    if comp_table:
        metadata["composition_table"] = comp_table
        if "elements" in comp_table and "values" in comp_table:
            for elem, val in zip(comp_table["elements"], comp_table["values"]):
                metadata["composition"][elem] = val

    # Build entity_type + entity_id (kept as-is for display)
    if metadata.get("nuance"):
        metadata["entity_type"] = "Nuance"
        metadata["entity_id"] = metadata["nuance"]
    elif metadata.get("product_name"):
        metadata["entity_type"] = "Matière première"
        ref = (metadata.get("reference") or "").strip()
        metadata["entity_id"] = f"{metadata['product_name']}|{ref}".strip("|")
    elif metadata.get("grade"):
        metadata["entity_type"] = "Grade"
        metadata["entity_id"] = str(metadata["grade"]).strip()
    else:
        metadata["entity_type"] = "Inconnu"
        metadata["entity_id"] = None

    return metadata


def extract_images_from_slide(slide, slide_number: int, output_dir: Path, ppt_name: str) -> List[str]:
    """Extract all images from a slide and save them."""
    image_paths = []
    img_count = 0

    for shape in slide.shapes:
        if hasattr(shape, "image"):
            try:
                image = shape.image
                image_bytes = image.blob

                filename = f"{ppt_name}_slide{slide_number:03d}_img{img_count:02d}.png"
                filepath = output_dir / filename

                img = Image.open(io.BytesIO(image_bytes))
                img.save(filepath, "PNG")

                image_paths.append(str(filepath.relative_to(output_dir.parent)))
                img_count += 1

            except Exception as e:
                print(f"   ⚠️  Error extracting image from slide {slide_number}: {e}")

    return image_paths


def process_powerpoint(ppt_path: Path, output_dir: Path) -> List[Dict]:
    """Process a PowerPoint file and extract all micrographs with metadata."""
    prs = Presentation(ppt_path)
    ppt_name = ppt_path.stem

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    print(f"\n📊 Processing: {ppt_path.name}")
    print(f"   Found {len(prs.slides)} slides")

    # First pass: extract metadata from ALL slides
    all_slide_metadata: List[Dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        metadata = extract_metadata_from_slide(slide, slide_idx)
        all_slide_metadata.append(metadata)

    # --- LOGIC: Entity Inheritance from NEXT slides ---
    # ✅ PATCH: propagate based on group_key (stable by reference for MPs)
    last_found_entity = None
    last_found_group = None

    for i in range(len(all_slide_metadata) - 1, -1, -1):
        sm = all_slide_metadata[i]
        gk = group_key(sm)

        if sm.get("entity_id") and gk:
            last_found_group = gk
            last_found_entity = {
                "entity_id": sm["entity_id"],
                "entity_type": sm["entity_type"],
                "nuance": sm.get("nuance"),
                "product_name": sm.get("product_name"),
                "reference": sm.get("reference"),
                "grade": sm.get("grade"),
            }
        elif last_found_entity:
            # inherit from subsequent slide
            sm["entity_id"] = last_found_entity["entity_id"]
            sm["entity_type"] = last_found_entity["entity_type"]
            sm["nuance"] = last_found_entity.get("nuance")
            sm["product_name"] = last_found_entity.get("product_name")
            sm["reference"] = last_found_entity.get("reference")
            sm["grade"] = last_found_entity.get("grade")
            print(f"   🔗 Slide {i+1} inherited entity '{sm['entity_id']}' from a subsequent slide")

    # Build aggregated entity_full_text per group key
    entity_text_map: Dict[str, str] = {}

    for sm in all_slide_metadata:
        gk = group_key(sm)
        if not gk:
            continue

        chunks = []
        ft = (sm.get("full_text") or "").strip()
        if ft:
            chunks.append(ft)

        comm = (sm.get("comments") or sm.get("description") or "").strip()
        if comm and comm not in ft:
            chunks.append(comm)

        ct = composition_table_to_text(sm.get("composition_table"))
        if ct:
            chunks.append("TABLEAU DE COMPOSITION:\n" + ct)

        if chunks:
            prev = entity_text_map.get(gk, "")
            merged = (prev + "\n\n" + "\n".join(chunks)).strip() if prev else "\n".join(chunks)
            entity_text_map[gk] = merged

    # Second pass: extract images and attach metadata
    all_metadata: List[Dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        current_metadata = all_slide_metadata[slide_idx - 1]
        image_paths = extract_images_from_slide(slide, slide_idx, images_dir, ppt_name)

        if image_paths:
            # Look for best comments (lookback up to 3 slides)
            best_comments = current_metadata.get("comments")

            # ✅ PATCH: compare by group_key, not raw entity_id
            current_gk = group_key(current_metadata)

            if not best_comments or len(best_comments) < 50:
                for lookback in range(1, 4):
                    prev_idx = slide_idx - 1 - lookback
                    if prev_idx >= 0:
                        prev_metadata = all_slide_metadata[prev_idx]
                        if group_key(prev_metadata) == current_gk:
                            prev_comm = prev_metadata.get("comments") or ""
                            if prev_comm and len(prev_comm) > 50:
                                best_comments = prev_comm
                                print(f"   📝 Using comments from slide {prev_idx + 1} for images in slide {slide_idx}")
                                break

            # Create entries for each image
            for img_path in image_paths:
                entry = current_metadata.copy()
                entry["image_path"] = img_path
                entry["source_file"] = ppt_path.name

                if best_comments:
                    entry["comments"] = best_comments
                    entry["description"] = best_comments

                gk = group_key(entry)
                if gk and gk in entity_text_map:
                    entry["entity_full_text"] = entity_text_map[gk]
                else:
                    entry["entity_full_text"] = ""

                all_metadata.append(entry)

                if entry.get("entity_type") and entry.get("entity_id"):
                    ref_info = f" [{entry['reference']}]" if entry.get("reference") else ""
                    desc_info = " (with description)" if best_comments and len(best_comments) > 50 else ""
                    print(f"   ✅ Slide {slide_idx}: {entry['entity_type']}={entry['entity_id']}{ref_info}{desc_info}")
                else:
                    print(f"   ✅ Slide {slide_idx}: {len(image_paths)} image(s)")

    return all_metadata


def main():
    parser = argparse.ArgumentParser(
        description="Extract micrographs and metadata from PowerPoint files (v3.6 patched)"
    )
    parser.add_argument("--input_dir", type=str, required=True, help="Directory containing PowerPoint files")
    parser.add_argument("--output_dir", type=str, required=True, help="Output directory for images and metadata")

    args = parser.parse_args()

    input_path = Path(args.input_dir)
    output_path = Path(args.output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    ppt_files = list(input_path.glob("*.pptx")) + list(input_path.glob("*.ppt"))

    if not ppt_files:
        print(f"❌ No PowerPoint files found in {input_path}")
        return

    print(f"📚 Found {len(ppt_files)} PowerPoint file(s)")

    all_metadata: List[Dict] = []
    for ppt_file in ppt_files:
        metadata = process_powerpoint(ppt_file, output_path)
        all_metadata.extend(metadata)

    metadata_path = output_path / "metadata.json"
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(all_metadata, f, indent=2, ensure_ascii=False)

    print(f"\n🎉 Extraction complete!")
    print(f"   📸 Total images extracted: {len(all_metadata)}")
    print(f"   📁 Images saved to: {output_path / 'images'}")
    print(f"   🧾 Metadata saved to: {metadata_path}")


if __name__ == "__main__":
    main()
