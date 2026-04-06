#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║        EXTRACTION MICROGRAPHIES — VERSION 4 (UNIFIÉE) — FIXED              ║
║                                                                              ║
║  FIX: All embeddings are now L2-normalized before DB insertion.             ║
║       compute_embedding_from_pil() → _normalize_embedding() → store        ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

import os
import re
import io
import shutil
import tempfile
import numpy as np
from pathlib import Path
from typing import Dict, List, Optional, Any
from urllib.parse import urlparse, unquote

import torch
from transformers import AutoModel, AutoImageProcessor
from pptx import Presentation
from PIL import Image
from azure.storage.blob import BlobServiceClient
from azure.core.exceptions import ResourceExistsError
from dotenv import load_dotenv

import psycopg2
from psycopg2.extras import Json
from pgvector.psycopg2 import register_vector
from pgvector import Vector

# ──────────────────────────────────────────────────────────────────────────────
# CONFIGURATION
# ──────────────────────────────────────────────────────────────────────────────
load_dotenv()

DB_DSN = os.getenv(
    "DB_DSN",
    "postgresql://administrationSTS:St%24%400987@avo-adb-002.postgres.database.azure.com:5432/Micrographie_IA",
).strip()
AZURE_CONNECTION_STRING = os.getenv("AZURE_CONNECTION_STRING", "").strip()
AZURE_CONTAINER_NAME = os.getenv("AZURE_CONTAINER_NAME", "").strip()
AZURE_BLOB_PREFIX = os.getenv("AZURE_BLOB_PREFIX", "micrograph-images").strip("/ ")
_BLOB_SERVICE_CLIENT: Optional[BlobServiceClient] = None

# ──────────────────────────────────────────────────────────────────────────────
# SETUP DINOv2
# ──────────────────────────────────────────────────────────────────────────────
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"🔧 Loading DINOv2-large on {device}...")
dinov2_model = AutoModel.from_pretrained("facebook/dinov2-large").to(device).eval()
dinov2_processor = AutoImageProcessor.from_pretrained("facebook/dinov2-large")
print("✅ DINOv2 model loaded\n")

BASE_DIR = Path(__file__).resolve().parent
AZURE_PPT_TARGET_FOLDER = os.getenv("AZURE_PPT_TARGET_FOLDER", "micrographie-ppts-inputs").strip("/ ")


def is_azure_blob_enabled() -> bool:
    return bool(AZURE_CONNECTION_STRING and AZURE_CONTAINER_NAME)


def get_blob_service_client() -> BlobServiceClient:
    global _BLOB_SERVICE_CLIENT
    if not is_azure_blob_enabled():
        raise RuntimeError(
            "Azure Blob Storage is not configured. Set AZURE_CONNECTION_STRING and AZURE_CONTAINER_NAME."
        )
    if _BLOB_SERVICE_CLIENT is None:
        _BLOB_SERVICE_CLIENT = BlobServiceClient.from_connection_string(AZURE_CONNECTION_STRING)
        container_client = _BLOB_SERVICE_CLIENT.get_container_client(AZURE_CONTAINER_NAME)
        try:
            container_client.create_container()
        except ResourceExistsError:
            pass
    return _BLOB_SERVICE_CLIENT


def _build_blob_name(output_dir: Path, filename: str) -> str:
    safe_name = Path(filename).name
    folder_name = output_dir.name or "images"
    if AZURE_BLOB_PREFIX:
        return f"{AZURE_BLOB_PREFIX}/{folder_name}/{safe_name}"
    return f"{folder_name}/{safe_name}"


def _blob_path_from_name(blob_name: str) -> str:
    return f"azure-blob://{AZURE_CONTAINER_NAME}/{blob_name}"


def parse_azure_blob_path(storage_path: str) -> Optional[tuple[str, str]]:
    raw_value = str(storage_path or "").strip()
    if not raw_value:
        return None

    if raw_value.startswith("azure-blob://"):
        remainder = raw_value[len("azure-blob://"):]
        container_name, _, blob_name = remainder.partition("/")
        if container_name and blob_name:
            return container_name, blob_name
        return None

    parsed = urlparse(raw_value)
    path = unquote(parsed.path or "").lstrip("/")
    if parsed.scheme in {"http", "https"} and parsed.netloc.endswith(".blob.core.windows.net"):
        container_name, _, blob_name = path.partition("/")
        if container_name and blob_name:
            return container_name, blob_name

    return None


def build_powerpoint_blob_storage_path(source_path: str) -> Optional[str]:
    normalized = (source_path or "").replace("\\", "/").strip()
    if not normalized:
        return None
    if normalized.startswith("azure-blob://"):
        return normalized

    filename = Path(normalized).name
    if not filename:
        return None

    blob_name = f"{AZURE_BLOB_PREFIX}/{AZURE_PPT_TARGET_FOLDER}/{filename}" if AZURE_BLOB_PREFIX else f"{AZURE_PPT_TARGET_FOLDER}/{filename}"
    return f"azure-blob://{AZURE_CONTAINER_NAME}/{blob_name}"


def materialize_powerpoint_path(file_path_value: str) -> tuple[Optional[Path], Optional[Path]]:
    azure_blob_location = parse_azure_blob_path(file_path_value)
    if azure_blob_location:
        container_name, blob_name = azure_blob_location
        temp_dir = Path(tempfile.mkdtemp(prefix="ppt_"))
        temp_path = temp_dir / Path(blob_name).name
        blob_client = get_blob_service_client().get_blob_client(container=container_name, blob=blob_name)
        with open(temp_path, "wb") as stream:
            stream.write(blob_client.download_blob().readall())
        return temp_path, temp_dir

    local_path = Path(file_path_value)
    if local_path.exists():
        return local_path, None

    return None, None


def store_pil_image(image: Image.Image, output_dir: Path, filename: str, *, format: str = "PNG") -> str:
    safe_name = Path(filename).name
    if is_azure_blob_enabled():
        buffer = io.BytesIO()
        image.save(buffer, format=format)
        blob_name = _build_blob_name(output_dir, safe_name)
        blob_client = get_blob_service_client().get_blob_client(
            container=AZURE_CONTAINER_NAME,
            blob=blob_name,
        )
        blob_client.upload_blob(
            buffer.getvalue(),
            overwrite=True,
            content_type="image/png" if format.upper() == "PNG" else "application/octet-stream",
        )
        return _blob_path_from_name(blob_name)

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    filepath = images_dir / safe_name
    image.save(filepath, format)
    try:
        return str(filepath.relative_to(output_dir.parent))
    except ValueError:
        return str(filepath)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — UTILITAIRES COMMUNS
# ══════════════════════════════════════════════════════════════════════════════

def clean_text(text: str) -> str:
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def clean_reference(ref: str) -> str:
    if not ref:
        return ""
    ref = re.sub(r'^(ref|référence)\s*[:\-]?\s*', '', ref, flags=re.IGNORECASE)
    ref = re.sub(r'\s+', '', ref)
    return ref.strip().upper()


def compute_embedding_from_pil(image: Image.Image) -> np.ndarray:
    """Calcule l'embedding DINOv2 brut (1024 dimensions, non normalisé)."""
    image = image.convert("RGB")
    inputs = dinov2_processor(images=image, return_tensors="pt")
    inputs = {k: v.to(device) for k, v in inputs.items()}
    with torch.no_grad():
        outputs = dinov2_model(**inputs)
        embedding = outputs.last_hidden_state[:, 0, :].squeeze().cpu().numpy()
    return embedding.astype("float32")


# ✅ FIX — New helper: always normalize before storing in DB
def _normalize_embedding(embedding: np.ndarray) -> np.ndarray:
    """L2-normalise un embedding pour que pgvector cosine distance soit correcte."""
    arr = np.asarray(embedding, dtype="float32").ravel()
    norm = float(np.linalg.norm(arr))
    if norm <= 0:
        raise ValueError("Zero-norm embedding cannot be normalized")
    return (arr / norm).astype("float32")


# ✅ FIX — Multi-scale query embedding (mirrors compute_query_embedding_from_pil in app.py)
def compute_normalized_embedding(image: Image.Image) -> np.ndarray:
    """
    Calcule un embedding L2-normalisé depuis 4 vues multi-échelles.
    MUST match compute_query_embedding_from_pil() in app.py exactly.
    Uses the same crop ratios: [full, 0.92, 0.80, 0.66].
    """
    image = image.convert("RGB")
    width, height = image.size

    views: List[Image.Image] = [image]
    seen = {(width, height, 0, 0)}
    for ratio in [0.92, 0.80, 0.66]:
        crop_w = max(32, int(width * ratio))
        crop_h = max(32, int(height * ratio))
        left   = max(0, (width - crop_w) // 2)
        top    = max(0, (height - crop_h) // 2)
        key    = (crop_w, crop_h, left, top)
        if key in seen:
            continue
        seen.add(key)
        views.append(image.crop((left, top, left + crop_w, top + crop_h)))

    view_embeddings = [_normalize_embedding(compute_embedding_from_pil(v)) for v in views]
    averaged = np.mean(np.stack(view_embeddings, axis=0), axis=0)
    return _normalize_embedding(averaged)


def extract_magnifications_with_positions(slide) -> List[Dict]:
    """Extrait tous les grossissements et leur position verticale."""
    mags = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            match = re.search(r'[Gg]rossissement\s*[xX×]\s*(\d+)', shape.text)
            if match:
                mags.append({
                    "value": int(match.group(1)),
                    "text": shape.text.strip(),
                    "top": shape.top
                })
    return sorted(mags, key=lambda x: x["top"])


def assign_magnification(shape_top: int, magnifications: List[Dict]) -> Optional[int]:
    """Retourne le grossissement le plus proche au-dessus d'une image."""
    if not magnifications:
        return None
    mags_above = [m for m in magnifications if m["top"] <= shape_top]
    if mags_above:
        return mags_above[-1]["value"]
    return magnifications[0]["value"]


def extract_tables_from_slide(slide) -> List[List[List[str]]]:
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            try:
                tbl = shape.table
                table_rows = []
                for r in tbl.rows:
                    row_cells = [(c.text or "").strip() for c in r.cells]
                    table_rows.append(row_cells)
                tables.append(table_rows)
            except Exception:
                continue
    return tables


def _extract_nuance_from_text(text: str) -> Optional[str]:
    if not text:
        return None
    t = " ".join(str(text).split())
    m = re.search(r"\bNuance\s*[:\-]?\s*(\d{3})\s*(\d{2})\b", t, flags=re.IGNORECASE)
    if m:
        return f"{m.group(1)} {m.group(2)}"
    m = re.search(r"\b(\d{3})\s+(\d{2})\b", t)
    if m:
        return f"{m.group(1)} {m.group(2)}"
    m = re.search(r"\b(\d{5})\b", t)
    if m:
        s = m.group(1)
        return f"{s[:3]} {s[3:]}"
    return None


def extract_detailed_comments(text: str) -> Optional[str]:
    if not text:
        return None
    comments_match = re.search(
        r"Commentaires?\s*:\s*(.+)", text, re.IGNORECASE | re.DOTALL
    )
    if comments_match:
        comments = comments_match.group(1).strip()
        for pattern in [r'\bTABLE\b', r'\bComposition\b', r'\bGrossissement\s+x\s*\d+']:
            match = re.search(pattern, comments, flags=re.IGNORECASE)
            if match:
                comments = comments[:match.start()].strip()
                break
        comments = re.sub(r'[ \t]+', ' ', comments)
        comments = re.sub(r'\n{3,}', '\n\n', comments)
        comments = comments.strip()
        if len(comments) > 20:
            return comments
    return None


def parse_avo_composition_from_tables(tables: List[List[List[str]]]) -> Optional[Dict[str, Any]]:
    """Extrait un tableau de composition (format AVO) depuis les tables PPT."""
    for t in tables:
        rows = [[(c or "").strip() for c in r] for r in t if any((c or "").strip() for c in r)]
        if len(rows) < 2:
            continue
        for idx in range(1, len(rows)):
            value_row = rows[idx]
            numeric_cells = [c for c in value_row if c and re.match(r"^[<>]?\s*\d+(?:[.,]\d+)?\s*$", c)]
            if len(numeric_cells) < 2:
                continue
            header_row = rows[0]
            pairs = []
            for j in range(min(len(header_row), len(value_row))):
                h = header_row[j]
                v = value_row[j]
                if h and v and re.match(r"^[<>]?\s*\d+(?:[.,]\d+)?\s*$", v):
                    pairs.append((h, v))
            if len(pairs) >= 2:
                return {
                    "elements": [h for h, _ in pairs],
                    "values": [v for _, v in pairs],
                    "rows": [[h for h, _ in pairs], [v for _, v in pairs]]
                }
    return None


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — DÉTECTION DU TYPE DE FICHIER
# ══════════════════════════════════════════════════════════════════════════════

def is_metallic_nuances_file(prs: Presentation) -> bool:
    slides = list(prs.slides)
    if len(slides) < 2:
        return False
    slide1_text = " ".join(
        shape.text for shape in slides[0].shapes if hasattr(shape, "text")
    ).lower()
    if "nuance" not in slide1_text and "métallique" not in slide1_text:
        return False
    slide2_text = " ".join(
        shape.text for shape in slides[1].shapes if hasattr(shape, "text")
    )
    slide2_has_images = any(hasattr(s, "image") for s in slides[1].shapes)
    return "Nuance" in slide2_text and "Commentaires" in slide2_text and not slide2_has_images


def is_cokes_comparative_file(prs: Presentation) -> bool:
    slides = list(prs.slides)
    if len(slides) < 2:
        return False
    slide1_text = " ".join(
        shape.text for shape in slides[0].shapes if hasattr(shape, "text")
    )
    if "Coke" not in slide1_text:
        return False
    slide2_text = "\n".join(
        shape.text for shape in slides[1].shapes if hasattr(shape, "text")
    )
    return slide2_text.count("Micrographie N°") >= 3


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — BASE DE DONNÉES : CIRCUIT NUANCES MÉTALLIQUES
# ══════════════════════════════════════════════════════════════════════════════

def get_or_create_nuance(cur, conn, code_nuance: str, reference: Optional[str]) -> Optional[int]:
    if not code_nuance:
        return None
    if reference:
        cur.execute("SELECT id FROM public.nuances WHERE reference = %s", (reference,))
    else:
        cur.execute("SELECT id FROM public.nuances WHERE name = %s", (code_nuance,))
    row = cur.fetchone()
    if row:
        return row[0]
    cur.execute(
        "INSERT INTO public.nuances (reference, name, status, created_at) VALUES (%s, %s, 'active', NOW()) RETURNING id",
        (reference or code_nuance, code_nuance),
    )
    new_id = cur.fetchone()[0]
    conn.commit()
    return new_id


def insert_nuance_image(
    conn, cur,
    nuance_id: int,
    image_path: str,
    embedding: np.ndarray,   # ← must be L2-normalized before calling
    code_nuance: str,
    reference: Optional[str],
    metadata: Dict,
    source_file_id: int,
) -> bool:
    try:
        embedding_vector = Vector(embedding.tolist())
        cur.execute(
            """
            INSERT INTO public.nuance_images
                (nuance_id, image_path, embedding, nuance_name, reference, created_at)
            VALUES (%s, %s, %s, %s, %s, NOW())
            RETURNING id
            """,
            (nuance_id, image_path, embedding_vector, code_nuance, reference),
        )
        image_id = cur.fetchone()[0]
        note_json = {
            "expert_notes":   metadata.get("comments") or "",
            "magnification":  metadata.get("magnification"),
            "slide_number":   metadata.get("slide_number"),
            "source_file_id": source_file_id,
            "nuance":         code_nuance,
            "reference":      reference,
            "composition":    metadata.get("composition"),
            "annotations":    metadata.get("annotations", []),
            "type":           "Nuance métallique",
        }
        cur.execute(
            "INSERT INTO public.nuance_expert_notes (nuance_image_id, note_json, created_at) VALUES (%s, %s, NOW())",
            (image_id, Json(note_json)),
        )
        conn.commit()
        return True
    except Exception as e:
        conn.rollback()
        print(f"      ❌ Erreur insertion nuance_images: {e}")
        return False


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — BASE DE DONNÉES : CIRCUIT MATIÈRES PREMIÈRES
# ══════════════════════════════════════════════════════════════════════════════

def get_or_create_matiere(cur, conn, entry: Dict) -> Optional[int]:
    ref = entry.get("reference")
    if not ref:
        return None
    cur.execute("SELECT matiere_id FROM public.matieres WHERE reference = %s", (ref,))
    row = cur.fetchone()
    if row:
        return row[0]
    name = entry.get("product_name") or f"Matière {ref}"
    matiere_type = entry.get("type_matiere") or "Matière première"
    cur.execute(
        """
        INSERT INTO public.matieres (nom_matiere, type_matiere, reference, date_creation, date_mise_a_jour)
        VALUES (%s, %s, %s, NOW(), NOW())
        RETURNING matiere_id
        """,
        (name, matiere_type, ref),
    )
    new_id = cur.fetchone()[0]
    conn.commit()
    return new_id


def insert_matiere_image(
    conn, cur,
    matiere_id: int,
    image_path: str,
    embedding: np.ndarray,   # ← must be L2-normalized before calling
    entry: Dict,
    source_file_id: int,
) -> bool:
    try:
        embedding_vector = Vector(embedding.tolist())
        cur.execute(
            """
            INSERT INTO public.matiere_images
                (matiere_id, image_path, embedding, material_name, reference)
            VALUES (%s, %s, %s, %s, %s)
            RETURNING id
            """,
            (
                matiere_id,
                image_path,
                embedding_vector,
                entry.get("product_name"),
                entry.get("reference"),
            ),
        )
        image_id = cur.fetchone()[0]
        note_json = {
            "expert_notes":   entry.get("comments") or "",
            "magnification":  entry.get("magnification"),
            "slide_number":   entry.get("slide_number"),
            "source_file_id": source_file_id,
            "reference":      entry.get("reference"),
            "composition":    entry.get("composition"),
            "type":           entry.get("type_matiere", "Matière première"),
        }
        cur.execute(
            "INSERT INTO public.matiere_expert_notes (matiere_image_id, note_json, created_at) VALUES (%s, %s, NOW())",
            (image_id, Json(note_json)),
        )
        conn.commit()
        return True
    except Exception as e:
        conn.rollback()
        print(f"      ❌ Erreur insertion matiere_images: {e}")
        return False


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — PROCESSEUR : NUANCES MÉTALLIQUES
# ══════════════════════════════════════════════════════════════════════════════

def _parse_metallic_composition_table(table) -> Optional[Dict[str, Any]]:
    rows = [[c.text.strip() for c in r.cells] for r in table.rows]
    rows = [r for r in rows if any(r)]
    if len(rows) < 2:
        return None
    headers = rows[0]
    value_row = None
    for r in reversed(rows):
        if any(re.match(r"^[\d,\.]+$", c) for c in r if c):
            value_row = r
            break
    if not value_row:
        return None
    ref_row = rows[1] if len(rows) > 2 else None
    components = []
    for j, header in enumerate(headers):
        if not header:
            continue
        val = value_row[j] if j < len(value_row) else ""
        ref = ref_row[j] if ref_row and j < len(ref_row) else ""
        if header and val:
            components.append({"name": header, "reference": ref or None, "percentage": val})
    return {"components": components} if components else None


def _parse_metallic_metadata_slide(slide) -> Dict[str, Any]:
    result: Dict[str, Any] = {
        "nuance": None, "reference": None, "comments": None, "composition": None
    }
    for shape in slide.shapes:
        text = shape.text.strip() if hasattr(shape, "text") else ""
        if "Nuance" in text and result["nuance"] is None:
            m = re.search(r"Nuance\s+([A-Z0-9]{3}\s+[A-Z0-9]{2,3})", text, re.IGNORECASE)
            if m:
                result["nuance"] = m.group(1).strip()
            ref_m = re.search(r"\(\s*(\d{4}\s+\d{3})\s*\)", text)
            if ref_m:
                result["reference"] = ref_m.group(1).replace(" ", "")
        if "Commentaires" in text and result["comments"] is None:
            m = re.search(r"Commentaires\s*:\s*\n?(.*)", text, re.DOTALL)
            if m:
                raw = re.sub(r"[ \t]+", " ", m.group(1).strip())
                raw = re.sub(r"\n{3,}", "\n\n", raw)
                if len(raw) > 10:
                    result["comments"] = raw
        if getattr(shape, "has_table", False) and result["composition"] is None:
            result["composition"] = _parse_metallic_composition_table(shape.table)
    return result


def _group_slides_by_nuance(slides: list) -> List[Dict]:
    groups = []
    current_group = None
    for i, slide in enumerate(slides):
        slide_text = " ".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )
        has_images = any(hasattr(s, "image") for s in slide.shapes)
        is_metadata = (
            "Nuance" in slide_text
            and "Commentaires" in slide_text
            and not has_images
        )
        if is_metadata:
            current_group = {
                "metadata_slide": slide,
                "image_slides": [],
                "slide_indices": {"metadata": i + 1, "images": []},
            }
            groups.append(current_group)
        elif current_group is not None and has_images:
            current_group["image_slides"].append(slide)
            current_group["slide_indices"]["images"].append(i + 1)
    return groups


def _extract_annotations_from_image_slide(slide) -> List[str]:
    annotations = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if (
                len(text) < 60
                and "Grossissement" not in text
                and not re.match(r"^\d+$", text)
                and not re.match(r"^\d{3}\s+[A-Z0-9]{2,3}$", text)
                and "Commentaires" not in text
            ):
                annotations.append(clean_text(text))
    return annotations


def _extract_slide_text_blocks(slide) -> List[str]:
    return [
        clean_text(shape.text)
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    ]


def _extract_primary_title(slide) -> Optional[str]:
    candidates = []
    for text in _extract_slide_text_blocks(slide):
        if "Grossissement" in text or re.fullmatch(r"[Xx]\s*\d+", text):
            continue
        candidates.append(text)
    if not candidates:
        return None
    return max(candidates, key=len)


def _normalize_lookup_key(value: Optional[str]) -> str:
    return re.sub(r"[^a-z0-9]+", "", (value or "").casefold())


def _parse_name_and_reference(title: Optional[str]) -> tuple[Optional[str], Optional[str]]:
    cleaned_title = clean_text(title or "")
    if not cleaned_title:
        return None, None

    match = re.search(r"^(.*?)\s*[–-]?\s*ref\s+([A-Z0-9 ]+)$", cleaned_title, re.IGNORECASE)
    if match:
        return match.group(1).strip(" -–:"), clean_reference(match.group(2))
    return cleaned_title.strip(" -–:"), None


def _parse_standard_metadata_slide(slide) -> Dict[str, Any]:
    title = _extract_primary_title(slide)
    product_name, reference = _parse_name_and_reference(title)
    return {
        "product_name": product_name,
        "reference": reference,
        "comments": extract_detailed_comments("\n".join(_extract_slide_text_blocks(slide))),
        "composition": parse_avo_composition_from_tables(extract_tables_from_slide(slide)),
        "type_matiere": "Matière première",
    }


def _parse_cokes_metadata(prs: Presentation) -> List[Dict[str, Any]]:
    slides = list(prs.slides)
    if len(slides) < 2:
        return []

    comparison_text = "\n".join(_extract_slide_text_blocks(slides[1]))
    if not comparison_text:
        return []

    pattern = re.compile(
        r"(?P<title>[^:]+?)\s*:\s*Micrographie\s*N°\s*\d+\s*(?P<comments>.*?)(?=(?:[^:]+?)\s*:\s*Micrographie\s*N°\s*\d+|$)",
        re.IGNORECASE | re.DOTALL,
    )
    entries: List[Dict[str, Any]] = []
    for match in pattern.finditer(comparison_text):
        raw_title = clean_text(match.group("title"))
        comments = clean_text(match.group("comments").replace("Micrographies Coke", ""))
        product_name, reference = _parse_name_and_reference(raw_title)
        entries.append(
            {
                "product_name": product_name,
                "reference": reference,
                "comments": comments or None,
                "type_matiere": "Coke",
                "lookup_key": _normalize_lookup_key(product_name or raw_title),
            }
        )
    return entries


def _resolve_cokes_entry(title: Optional[str], entries: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    title_key = _normalize_lookup_key(title)
    if not title_key:
        return None

    exact_match = next((entry for entry in entries if entry["lookup_key"] == title_key), None)
    if exact_match:
        return exact_match

    partial_match = next(
        (
            entry
            for entry in entries
            if entry["lookup_key"] in title_key or title_key in entry["lookup_key"]
        ),
        None,
    )
    return partial_match


def _process_material_image_shapes(
    conn,
    cur,
    ppt_path: Path,
    output_dir: Path,
    file_id: int,
    slide,
    slide_num: int,
    entry: Dict[str, Any],
    img_start_index: int,
) -> int:
    magnifications = extract_magnifications_with_positions(slide)
    annotations = _extract_annotations_from_image_slide(slide)
    inserted = 0
    image_index = img_start_index

    reference = clean_reference(entry.get("reference") or "") or None
    material_entry = {
        "product_name": entry.get("product_name") or f"Matière {ppt_path.stem}",
        "reference": reference,
        "comments": entry.get("comments"),
        "slide_number": slide_num,
        "composition": entry.get("composition"),
        "type_matiere": entry.get("type_matiere") or "Matière première",
    }

    matiere_id = get_or_create_matiere(cur, conn, material_entry)
    if not matiere_id:
        print(f"      ⚠️  Matière ignorée faute de référence: {material_entry['product_name']}")
        return 0

    for shape in slide.shapes:
        if not hasattr(shape, "image"):
            continue

        mag = assign_magnification(shape.top, magnifications)
        img = Image.open(io.BytesIO(shape.image.blob))
        embedding = compute_normalized_embedding(img)
        filename = f"{ppt_path.stem}_s{slide_num:03d}_i{image_index:02d}.png"
        storage_path = store_pil_image(img, output_dir, filename, format="PNG")

        success = insert_matiere_image(
            conn,
            cur,
            matiere_id=matiere_id,
            image_path=storage_path,
            embedding=embedding,
            entry={
                **material_entry,
                "magnification": mag,
                "annotations": annotations,
            },
            source_file_id=file_id,
        )
        if success:
            inserted += 1
            mag_str = f"x{mag}" if mag else "?"
            print(
                f"      ✅ Slide {slide_num} img {image_index:02d} | "
                f"{material_entry['product_name']} | Gross. {mag_str}"
            )
        image_index += 1

    return inserted


def process_standard_powerpoint(ppt_path: Path, output_dir: Path, file_id: int):
    print(f"\n🧪 Traitement [MATIÈRES PREMIÈRES] : {ppt_path.name}")

    conn = psycopg2.connect(DB_DSN)
    register_vector(conn)
    cur = conn.cursor()

    try:
        prs = Presentation(ppt_path)
        current_entry: Optional[Dict[str, Any]] = None
        total_images = 0
        img_global = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            has_images = any(hasattr(shape, "image") for shape in slide.shapes)
            if not has_images:
                parsed_entry = _parse_standard_metadata_slide(slide)
                if parsed_entry.get("product_name"):
                    current_entry = parsed_entry
                continue

            if current_entry is None:
                fallback_name, fallback_ref = _parse_name_and_reference(_extract_primary_title(slide))
                current_entry = {
                    "product_name": fallback_name,
                    "reference": fallback_ref,
                    "comments": None,
                    "composition": None,
                    "type_matiere": "Matière première",
                }

            inserted = _process_material_image_shapes(
                conn,
                cur,
                ppt_path,
                output_dir,
                file_id,
                slide,
                slide_num,
                current_entry,
                img_global,
            )
            total_images += inserted
            img_global += sum(1 for shape in slide.shapes if hasattr(shape, "image"))

        print(f"\n   ✅ Terminé : {total_images} image(s) insérée(s) [matiere_images]")

    except Exception as e:
        print(f"❌ Erreur : {e}")
        import traceback; traceback.print_exc()
    finally:
        cur.close()
        conn.close()


def process_cokes_powerpoint(ppt_path: Path, output_dir: Path, file_id: int):
    print(f"\n🔥 Traitement [COKES COMPARATIFS] : {ppt_path.name}")

    conn = psycopg2.connect(DB_DSN)
    register_vector(conn)
    cur = conn.cursor()

    try:
        prs = Presentation(ppt_path)
        cokes_entries = _parse_cokes_metadata(prs)
        total_images = 0
        img_global = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            has_images = any(hasattr(shape, "image") for shape in slide.shapes)
            if not has_images:
                continue

            slide_title = _extract_primary_title(slide)
            entry = _resolve_cokes_entry(slide_title, cokes_entries)
            if entry is None:
                product_name, reference = _parse_name_and_reference(slide_title)
                entry = {
                    "product_name": product_name,
                    "reference": reference,
                    "comments": None,
                    "composition": None,
                    "type_matiere": "Coke",
                }

            inserted = _process_material_image_shapes(
                conn,
                cur,
                ppt_path,
                output_dir,
                file_id,
                slide,
                slide_num,
                entry,
                img_global,
            )
            total_images += inserted
            img_global += sum(1 for shape in slide.shapes if hasattr(shape, "image"))

        print(f"\n   ✅ Terminé : {total_images} image(s) insérée(s) [matiere_images]")

    except Exception as e:
        print(f"❌ Erreur : {e}")
        import traceback; traceback.print_exc()
    finally:
        cur.close()
        conn.close()


def process_metallic_powerpoint(ppt_path: Path, output_dir: Path, file_id: int):
    """Traite un fichier PowerPoint de nuances métalliques."""
    print(f"\n🔬 Traitement [NUANCES MÉTALLIQUES] : {ppt_path.name}")

    conn = psycopg2.connect(DB_DSN)
    register_vector(conn)
    cur = conn.cursor()

    try:
        prs = Presentation(ppt_path)
        slides = list(prs.slides)
        groups = _group_slides_by_nuance(slides)
        print(f"   → {len(groups)} nuances détectées")

        total_images = 0

        for group in groups:
            meta = _parse_metallic_metadata_slide(group["metadata_slide"])
            code_nuance = meta["nuance"]
            reference   = meta["reference"]
            comments    = meta["comments"]
            composition = meta["composition"]

            if not code_nuance:
                print(f"   ⚠️  Nuance non détectée dans la slide de métadonnées, ignorée")
                continue

            nuance_id = get_or_create_nuance(cur, conn, code_nuance, reference)
            if not nuance_id:
                print(f"   ⚠️  Impossible de créer la nuance {code_nuance}")
                continue

            ref_str = f"ref {reference}" if reference else "sans ref"
            print(f"\n   📌 Nuance {code_nuance} ({ref_str})")

            img_global = 0
            for slide_obj, slide_num in zip(
                group["image_slides"], group["slide_indices"]["images"]
            ):
                magnifications = extract_magnifications_with_positions(slide_obj)
                annotations    = _extract_annotations_from_image_slide(slide_obj)

                inline_comments = None
                for shape in slide_obj.shapes:
                    if hasattr(shape, "text") and "Commentaires" in shape.text:
                        m = re.search(r"Commentaires\s*:\s*\n?(.*)", shape.text, re.DOTALL)
                        if m:
                            raw = re.sub(r"[ \t]+", " ", m.group(1).strip())
                            if len(raw) > 10:
                                inline_comments = raw
                                break

                effective_comments = comments
                if inline_comments and not comments:
                    effective_comments = inline_comments
                elif inline_comments and comments:
                    effective_comments = comments + "\n\n[Annotations]\n" + inline_comments

                for shape in slide_obj.shapes:
                    if not hasattr(shape, "image"):
                        continue

                    mag = assign_magnification(shape.top, magnifications)
                    img_bytes = shape.image.blob
                    filename  = f"{ppt_path.stem}_s{slide_num:03d}_i{img_global:02d}.png"

                    img = Image.open(io.BytesIO(img_bytes))

                    # ✅ FIX: use normalized multi-scale embedding
                    embedding = compute_normalized_embedding(img)

                    storage_path = store_pil_image(img, output_dir, filename, format="PNG")

                    success = insert_nuance_image(
                        conn, cur,
                        nuance_id=nuance_id,
                        image_path=storage_path,
                        embedding=embedding,
                        code_nuance=code_nuance,
                        reference=reference,
                        metadata={
                            "comments":      effective_comments,
                            "composition":   composition,
                            "magnification": mag,
                            "slide_number":  slide_num,
                            "annotations":   annotations,
                        },
                        source_file_id=file_id,
                    )

                    if success:
                        mag_str = f"x{mag}" if mag else "?"
                        com_str = f"✓ ({len(effective_comments)} chars)" if effective_comments else "✗"
                        print(f"      ✅ Slide {slide_num} img {img_global:02d} | Gross. {mag_str} | Com. {com_str}")
                        total_images += 1
                    img_global += 1

        print(f"\n   ✅ Terminé : {total_images} image(s) insérée(s) [nuance_images]")

    except Exception as e:
        print(f"❌ Erreur : {e}")
        import traceback; traceback.print_exc()
    finally:
        cur.close()
        conn.close()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — PROCESSEUR : STANDARD (graphite, matières premières)
# ══════════════════════════════════════════════════════════════════════════════
# NOTE: process_cokes_powerpoint follows the same pattern.
# Apply the same fix everywhere: replace compute_embedding_from_pil(img)
# with compute_normalized_embedding(img) before any insert call.
#
# Search for ALL occurrences of:
#   embedding = compute_embedding_from_pil(img)
# and replace with:
#   embedding = compute_normalized_embedding(img)
#
# There are exactly 2 occurrences in the original file:
#   - Line 684  (process_metallic_powerpoint) ← fixed above
#   - Line 1064 (process_standard_powerpoint) ← fix the same way
# ══════════════════════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 8 — ROUTER PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════

def process_powerpoint(ppt_path: Path, output_dir: Path, file_id: int):
    try:
        prs = Presentation(ppt_path)
        if is_metallic_nuances_file(prs):
            process_metallic_powerpoint(ppt_path, output_dir, file_id)
        elif is_cokes_comparative_file(prs):
            process_cokes_powerpoint(ppt_path, output_dir, file_id)
        else:
            process_standard_powerpoint(ppt_path, output_dir, file_id)
    except Exception as e:
        print(f"❌ Erreur détection type {ppt_path.name}: {e}")
        import traceback; traceback.print_exc()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 9 — NETTOYAGE
# ══════════════════════════════════════════════════════════════════════════════

def clear_all_data(output_dir: Path):
    try:
        conn = psycopg2.connect(DB_DSN)
        cur = conn.cursor()
        print("🗑️  Nettoyage des données existantes...")
        cur.execute("DELETE FROM public.matiere_expert_notes")
        cur.execute("DELETE FROM public.matiere_images")
        print("   ✅ matiere_expert_notes / matiere_images vidées")
        cur.execute("DELETE FROM public.nuance_expert_notes")
        cur.execute("DELETE FROM public.nuance_images")
        print("   ✅ nuance_expert_notes / nuance_images vidées")
        conn.commit()
        cur.close()
        conn.close()
        if is_azure_blob_enabled():
            print("   Azure Blob active: remote blobs are not deleted by this cleanup")
        images_dir = output_dir / "images"
        if images_dir.exists():
            try:
                shutil.rmtree(images_dir)
                print("   ✅ Fichiers images supprimés")
            except PermissionError:
                print("   ⚠️  Dossier images verrouillé, continuation...")
        print("✅ Nettoyage terminé\n")
    except Exception as e:
        print(f"❌ Erreur nettoyage : {e}")
        import traceback; traceback.print_exc()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 10 — MAIN
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    output = BASE_DIR / "output_v4"
    output.mkdir(exist_ok=True)

    clear_all_data(output)

    try:
        conn = psycopg2.connect(DB_DSN)
        cur = conn.cursor()
        cur.execute("SELECT id, file_path FROM public.powerpoint_files ORDER BY id")
        rows = cur.fetchall()
        cur.close()
        conn.close()

        if rows:
            print(f"📂 {len(rows)} fichier(s) trouvé(s) en base\n")
            for row in rows:
                file_id  = row[0]
                stored_path = row[1]
                ppt_path, temp_dir = materialize_powerpoint_path(stored_path)
                try:
                    if ppt_path is not None and ppt_path.exists():
                        process_powerpoint(ppt_path, output, file_id)
                    else:
                        expected_azure_path = build_powerpoint_blob_storage_path(stored_path)
                        if expected_azure_path and expected_azure_path != stored_path:
                            print(
                                f"⚠️  Fichier introuvable : {stored_path} | "
                                f"Path Azure attendu: {expected_azure_path}"
                            )
                        else:
                            print(f"⚠️  Fichier introuvable : {stored_path}")
                finally:
                    if temp_dir is not None:
                        shutil.rmtree(temp_dir, ignore_errors=True)
        else:
            print("⚠️  Aucun fichier PowerPoint en base de données")

    except Exception as e:
        print(f"❌ Erreur lecture base : {e}")
        import traceback; traceback.print_exc()